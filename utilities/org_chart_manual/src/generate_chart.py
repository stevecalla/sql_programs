# org_chart_with_lines.py
import pandas as pd
from collections import defaultdict, deque
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INPUT = "org_chart_usat_test.csv"
OUTPUT = "org_chart_usat_test.pptx"

df = pd.read_csv(INPUT, dtype=str).fillna("")
need = ["employee_id","name","title","manager_id"]
for c in need:
    if c not in df.columns:
        raise ValueError(f"Missing column: {c}")
    df[c] = df[c].astype(str).str.strip()
df = df.drop_duplicates(subset=["employee_id"])

by_id = {r.employee_id: r for r in df.itertuples(index=False)}
children = defaultdict(list)
for r in df.itertuples(index=False):
    if r.manager_id and r.manager_id in by_id and r.manager_id != r.employee_id:
        children[r.manager_id].append(r.employee_id)

roots = [
    r.employee_id for r in df.itertuples(index=False)
    if (r.manager_id == "" or r.manager_id not in by_id or r.manager_id == r.employee_id)
]
if not roots:
    all_children = {c for subs in children.values() for c in subs}
    roots = [eid for eid in by_id if eid not in all_children] or [next(iter(by_id))]

# Build levels
level = {}
for rt in roots:
    q = deque([(rt, 0)])
    while q:
        u, d = q.popleft()
        if u in level and level[u] <= d:
            continue
        level[u] = d
        for v in children[u]:
            q.append((v, d + 1))

# Stable DFS ordering within each level
ordered_by_level = defaultdict(list)
def order_subtree(u):
    ordered_by_level[level[u]].append(u)
    for v in children[u]:
        order_subtree(v)
for r in roots:
    order_subtree(r)

# Create slide
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

left_margin, top_margin = Inches(0.5), Inches(0.6)
x_gap, y_gap = Inches(2.6), Inches(1.6)
box_w, box_h = Inches(2.4), Inches(1.0)

page_width = prs.slide_width - 2 * left_margin

def x_for(col, cols):
    total = cols * x_gap
    start = left_margin + max((page_width - total) / 2, 0)
    return start + col * x_gap

# Draw nodes
centers = {}
for lvl in sorted(ordered_by_level):
    row = ordered_by_level[lvl]
    for col, emp_id in enumerate(row):
        r = by_id[emp_id]
        x = x_for(col, len(row))
        y = top_margin + lvl * y_gap

        box = slide.shapes.add_shape(1, x, y, box_w, box_h)  # 1 = rectangle
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        box.line.color.rgb = RGBColor(120, 120, 120)

        tf = box.text_frame
        tf.clear()  # reset content
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # center vertically

        p1 = tf.paragraphs[0]
        p1.text = r.name
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0, 0, 0)
        p1.alignment = PP_ALIGN.CENTER  # center horizontally

        p2 = tf.add_paragraph()
        p2.text = r.title
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0, 0, 0)
        p2.alignment = PP_ALIGN.CENTER

        centers[emp_id] = (x + box_w/2, y + box_h/2)

# Helper to draw a thin rectangle line
def draw_line(x, y, w, h, color=(100,100,100)):
    if w < 0:
        x = x + w; w = -w
    if h < 0:
        y = y + h; h = -h
    line = slide.shapes.add_shape(1, x, y, max(Emu(1), w), max(Emu(1), h))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*color)
    line.line.fill.background()

# Draw lines
thin = Emu(3)
for mng, subs in children.items():
    if mng not in centers:
        continue
    mx, my = centers[mng]
    m_bottom_y = my + box_h/2

    for s in subs:
        if s not in centers:
            continue
        sx, sy = centers[s]
        s_top_y = sy - box_h/2
        mid_y = (m_bottom_y + s_top_y) / 2

        if abs(mid_y - m_bottom_y) > 0:
            draw_line(mx - thin/2, m_bottom_y, thin, mid_y - m_bottom_y)
        if abs(sx - mx) > 0:
            draw_line(min(mx, sx), mid_y - thin/2, abs(sx - mx), thin)
        if abs(s_top_y - mid_y) > 0:
            draw_line(sx - thin/2, mid_y, thin, s_top_y - mid_y)

prs.save(OUTPUT)
print(f"Saved {OUTPUT}")
