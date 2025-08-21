# app.py
import streamlit as st

# ---- Backward-compat shim for older Streamlit versions ----
try:
    DATA_EDITOR = st.data_editor
except AttributeError:
    DATA_EDITOR = st.experimental_data_editor  # fallback for Streamlit < 1.19

import re
from typing import Optional
import pandas as pd
from collections import defaultdict, deque
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ---------------------------
# UI CONFIG
# ---------------------------
st.set_page_config(page_title="Org Chart → PowerPoint", layout="centered")
st.title("Org Chart → PowerPoint (One-Pager)")
st.write(
    "Upload **CSV or Excel (.xlsx/.xls)** with any headers. "
    "Then (optionally) map them to: **employee_id, name, title, manager_id, department** "
    "(optional: **tenure**, **tenure_calc**)."
)

# ---------------------------
# Feature Toggles
# ---------------------------
# Set to False to hide the mapping UI entirely.
ENABLE_MAPPING_UI = True
# If mapping UI is hidden, choose whether to auto-map (True) or require canonical headers (False).
AUTO_MAP_WHEN_DISABLED = True

# ---------------------------
# Drawing / Layout constants (tunable)
# ---------------------------
LEFT_M_IN, RIGHT_M_IN = 0.5, 0.5
TOP_M_IN, BOTTOM_M_IN = 0.7, 0.4
TITLE_H_IN = 0.5

TARGET_BOX_W_IN = 2.4
TARGET_BOX_H_IN = 1.0
TARGET_X_GAP_IN = 2.6
TARGET_Y_GAP_IN = 1.4

MIN_BOX_W_IN = 1.2
MIN_BOX_H_IN = 0.75
MIN_X_GAP_IN = 0.35
MIN_Y_GAP_IN = 0.4

NAME_FONT_PT = 12
TITLE_FONT_PT = 10
TENURE_FONT_PT = 9
NAME_FONT_MIN_PT = 9
TITLE_FONT_MIN_PT = 8
TENURE_FONT_MIN_PT = 8

FOOTER_PT = 10
FOOTER_COLOR = (90, 90, 90)

TITLE_BAR_FILL = (245, 245, 245)
TITLE_BAR_LINE = (200, 200, 200)

BOX_FILL = (240, 240, 240)
BOX_LINE = (120, 120, 120)

ORDER_COL = "order"

REQUIRED_CANON = ["employee_id", "name", "title", "manager_id", "department"]
OPTIONAL_CANON = ["tenure", "tenure_calc"]
ALL_CANON = REQUIRED_CANON + OPTIONAL_CANON

SYNONYMS = {
    "employee_id": ["employee_id","employee id","emp id","emplid","eid","user id","user_id","uid","id","person id","personid"],
    "name":        ["name","full name","employee name","staff name"],
    "title":       ["title","job title","position","role"],
    "manager_id":  ["manager_id","manager id","mgr id","manager","reports to id","reports_to","supervisor id","supervisor"],
    "department":  ["department","dept","org","organization","team","division","group"],
    "tenure":      ["tenure","service","service time","tenure text","tenure (text)"],
    "tenure_calc": ["tenure_calc","tenure (yrs)","tenure years","years","years_of_service","yrs","years service","years in role"]
}

# ---------------------------
# Helpers (drawing)
# ---------------------------

def add_footer_pagenum(slide, prs, page_num, total_pages):
    txt = f"Page {page_num} of {total_pages}"
    w = Inches(2.2); h = Inches(0.3)
    x = prs.slide_width - Inches(RIGHT_M_IN) - w
    y = prs.slide_height - Inches(BOTTOM_M_IN) - h
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = txt
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(FOOTER_PT)
    p.font.color.rgb = RGBColor(*FOOTER_COLOR)

def draw_person_box(slide, x, y, w, h, name, title, tenure, name_pt, title_pt, tenure_pt,
                    fill_rgb=BOX_FILL, line_rgb=BOX_LINE):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill_rgb)
    box.line.color.rgb = RGBColor(*line_rgb)
    tf = box.text_frame; tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.text = name or ""
    p1.font.size = Pt(name_pt); p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title or ""
    p2.font.size = Pt(title_pt)
    p2.font.color.rgb = RGBColor(0, 0, 0)
    p2.alignment = PP_ALIGN.CENTER

    t = (tenure or "").strip()
    if t:
        p3 = tf.add_paragraph()
        p3.text = t
        p3.font.size = Pt(tenure_pt)
        p3.font.color.rgb = RGBColor(60, 60, 60)
        p3.alignment = PP_ALIGN.CENTER

    return box

def draw_line(slide, x, y, w, h, color=(100,100,100)):
    if w < 0: x = x + w; w = -w
    if h < 0: y = y + h; h = -h
    seg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, max(Emu(1), w), max(Emu(1), h))
    seg.fill.solid(); seg.fill.fore_color.rgb = RGBColor(*color)
    seg.line.fill.background()

def compute_layout(prs, max_cols, rows):
    slide_w, slide_h = prs.slide_width, prs.slide_height
    LEFT_M, RIGHT_M = Inches(LEFT_M_IN), Inches(RIGHT_M_IN)
    TOP_M, BOTTOM_M = Inches(TOP_M_IN), Inches(BOTTOM_M_IN)
    TITLE_H = Inches(TITLE_H_IN)
    avail_w = slide_w - LEFT_M - RIGHT_M
    avail_h = slide_h - TOP_M - BOTTOM_M - TITLE_H

    box_w = Inches(TARGET_BOX_W_IN); box_h = Inches(TARGET_BOX_H_IN)
    x_gap = Inches(TARGET_X_GAP_IN); y_gap = Inches(TARGET_Y_GAP_IN)
    name_pt, title_pt, tenure_pt = NAME_FONT_PT, TITLE_FONT_PT, TENURE_FONT_PT

    def total_w(cols, bw, xg): return cols*bw + (cols-1)*xg if cols > 0 else 0
    def total_h(rws, bh, yg): return rws*bh + (rws-1)*yg if rws > 0 else 0

    if max_cols > 0 and total_w(max_cols, box_w, x_gap) > avail_w:
        scale_w = avail_w / total_w(max_cols, box_w, x_gap)
        box_w *= scale_w; x_gap *= scale_w

    min_box_w = Inches(MIN_BOX_W_IN); min_x_gap = Inches(MIN_X_GAP_IN)
    if box_w < min_box_w:
        box_w = min_box_w
        if max_cols > 1:
            over_w = total_w(max_cols, box_w, x_gap) - avail_w
            if over_w > 0:
                x_gap = max(min_x_gap, x_gap - over_w/(max_cols-1))
        name_pt = max(NAME_FONT_MIN_PT, NAME_FONT_PT - 2)
        title_pt = max(TITLE_FONT_MIN_PT, TITLE_FONT_PT - 2)
        tenure_pt = max(TENURE_FONT_MIN_PT, TENURE_FONT_PT - 1)

    if rows > 0 and total_h(rows, box_h, y_gap) > avail_h:
        scale_h = avail_h / total_h(rows, box_h, y_gap)
        box_h *= scale_h; y_gap *= scale_h

    min_box_h = Inches(MIN_BOX_H_IN); min_y_gap = Inches(MIN_Y_GAP_IN)
    if box_h < min_box_h:
        box_h = min_box_h
        if rows > 1:
            over_h = total_h(rows, box_h, y_gap) - avail_h
            if over_h > 0:
                y_gap = max(min_y_gap, y_gap - over_h/(rows-1))
        name_pt = max(NAME_FONT_MIN_PT, name_pt - 1)
        title_pt = max(TITLE_FONT_MIN_PT, title_pt - 1)
        tenure_pt = max(TENURE_FONT_MIN_PT, tenure_pt - 1)

    max_cols_fit = 1
    bw_min, xg_min = min_box_w, min_x_gap
    if (bw_min + xg_min) > 0:
        max_cols_fit = max(1, int((avail_w + xg_min) // (bw_min + xg_min)))

    return box_w, box_h, x_gap, y_gap, name_pt, title_pt, tenure_pt, max_cols_fit

def draw_cover_slide(prs, dept_ordered, dept_counts, total_employees, page_num, total_pages, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    LEFT_M, RIGHT_M = Inches(LEFT_M_IN), Inches(RIGHT_M_IN)
    TOP_M, BOTTOM_M = Inches(TOP_M_IN), Inches(BOTTOM_M_IN)
    TITLE_H = Inches(TITLE_H_IN)

    title_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, LEFT_M, Inches(0.2),
                                         prs.slide_width - LEFT_M - RIGHT_M, TITLE_H)
    title_shape.fill.solid(); title_shape.fill.fore_color.rgb = RGBColor(*TITLE_BAR_FILL)
    title_shape.line.color.rgb = RGBColor(*TITLE_BAR_LINE)
    tf = title_shape.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = "Organization Overview"
    p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)

    subtitle = slide.shapes.add_textbox(LEFT_M, TOP_M + TITLE_H + Inches(0.2),
                                        prs.slide_width - LEFT_M - RIGHT_M, Inches(0.4))
    stf = subtitle.text_frame; stf.clear()
    sp = stf.paragraphs[0]
    sp.text = f"Total Employees: {total_employees}"
    sp.alignment = PP_ALIGN.CENTER; sp.font.size = Pt(16); sp.font.bold = True

    if "tenure_calc" in df.columns:
        tmp = df.copy()
        tmp["tenure_calc"] = pd.to_numeric(tmp["tenure_calc"], errors="coerce")
        dept_avg_tenure = tmp.groupby("department")["tenure_calc"].mean().round(1).to_dict()
    else:
        dept_avg_tenure = {}

    rows = len(dept_ordered) + 1
    cols = 3
    table_w = prs.slide_width - LEFT_M - RIGHT_M
    table_h = Inches(0.38) * rows
    table_x = LEFT_M
    table_y = TOP_M + TITLE_H + Inches(0.9)
    shape = slide.shapes.add_table(rows, cols, table_x, table_y, table_w, table_h)
    table = shape.table

    table.cell(0,0).text = "Department"
    table.cell(0,1).text = "Employees"
    table.cell(0,2).text = "Avg Tenure (yrs)"
    for j in range(cols):
        cell = table.cell(0, j)
        cell.text_frame.paragraphs[0].font.bold = True

    for i, dept in enumerate(dept_ordered, start=1):
        cnt = dept_counts.get(dept, 0)
        avg = dept_avg_tenure.get(dept, float("nan"))
        avg_txt = "" if pd.isna(avg) else str(avg)
        table.cell(i, 0).text = dept
        table.cell(i, 1).text = str(cnt)
        table.cell(i, 2).text = avg_txt

    add_footer_pagenum(slide, prs, page_num, total_pages)

def get_department_order(df):
    depts = df["department"].dropna().unique().tolist()
    if ORDER_COL in df.columns:
        tmp = df[["department", ORDER_COL]].copy()
        tmp["_ord"] = pd.to_numeric(tmp[ORDER_COL], errors="coerce")
        first_order = (tmp.dropna(subset=["_ord"]).groupby("department")["_ord"].first().to_dict())
    else:
        first_order = {}
    def key_func(d):
        if d in first_order:
            return (0, float(first_order[d]), d.lower())
        return (1, float("inf"), d.lower())
    return sorted(depts, key=key_func)

def draw_department_slide(prs, dept_name, sub_df, page_num, total_pages, by_id_all):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    LEFT_M, RIGHT_M = Inches(LEFT_M_IN), Inches(RIGHT_M_IN)
    TOP_M, BOTTOM_M = Inches(TOP_M_IN), Inches(BOTTOM_M_IN)
    TITLE_H = Inches(TITLE_H_IN)

    dept_count = len(sub_df)
    title_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, LEFT_M, Inches(0.2),
                                         prs.slide_width - LEFT_M - RIGHT_M, TITLE_H)
    title_shape.fill.solid(); title_shape.fill.fore_color.rgb = RGBColor(*TITLE_BAR_FILL)
    title_shape.line.color.rgb = RGBColor(*TITLE_BAR_LINE)
    tf = title_shape.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Department: {dept_name}  •  Employees: {dept_count}"
    p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)

    sub_df = sub_df.copy()
    by_id = {r.employee_id: r for r in sub_df.itertuples(index=False)}

    children = defaultdict(list)
    for r in sub_df.itertuples(index=False):
        if r.manager_id and r.manager_id in by_id and r.manager_id != r.employee_id:
            children[r.manager_id].append(r.employee_id)

    dept_roots = [
        r.employee_id for r in sub_df.itertuples(index=False)
        if (r.manager_id == "" or r.manager_id == r.employee_id or r.manager_id not in by_id)
    ]
    if not dept_roots:
        all_children = {c for subs in children.values() for c in subs}
        dept_roots = [eid for eid in by_id if eid not in all_children] or [next(iter(by_id))]

    ext_mgrs_order = []
    mgr_to_dept_roots = defaultdict(list)
    for r in sub_df.itertuples(index=False):
        if (r.employee_id in dept_roots and r.manager_id and
                r.manager_id not in by_id and r.manager_id in by_id_all and
                r.manager_id != r.employee_id):
            if r.manager_id not in ext_mgrs_order:
                ext_mgrs_order.append(r.manager_id)
            mgr_to_dept_roots[r.manager_id].append(r.employee_id)

    level = {}
    for rt in dept_roots:
        q = deque([(rt, 0)])
        while q:
            u, d = q.popleft()
            if u in level and level[u] <= d:
                continue
            level[u] = d
            for v in children[u]:
                q.append((v, d + 1))

    ordered_by_level = defaultdict(list)
    def order_subtree(u):
        ordered_by_level[level[u]].append(u)
        for v in children[u]:
            order_subtree(v)
    for r in dept_roots:
        order_subtree(r)

    tentative_rows = (1 if ext_mgrs_order else 0) + (1 + (max(level.values()) if level else 0))
    per_level_counts = [len(ordered_by_level[l]) for l in sorted(ordered_by_level)]
    tentative_max_cols = max(per_level_counts + [len(ext_mgrs_order or [])] or [1])

    box_w, box_h, x_gap, y_gap, name_pt, title_pt, tenure_pt, max_cols_fit = compute_layout(
        prs, tentative_max_cols, tentative_rows
    )

    layout_rows = []
    if ext_mgrs_order:
        if len(ext_mgrs_order) > max_cols_fit:
            for i in range(0, len(ext_mgrs_order), max_cols_fit):
                layout_rows.append((True, ext_mgrs_order[i:i+max_cols_fit]))
        else:
            layout_rows.append((True, ext_mgrs_order))

    for lvl in sorted(ordered_by_level):
        row_nodes = ordered_by_level[lvl]
        if len(row_nodes) > max_cols_fit:
            for i in range(0, len(row_nodes), max_cols_fit):
                layout_rows.append((False, row_nodes[i:i+max_cols_fit]))
        else:
            layout_rows.append((False, row_nodes))

    rows = len(layout_rows)
    max_cols = max((len(rnodes) for _, rnodes in layout_rows), default=1)
    box_w, box_h, x_gap, y_gap, name_pt, title_pt, tenure_pt, _ = compute_layout(prs, max_cols, rows)

    page_width = prs.slide_width - Inches(LEFT_M_IN) - Inches(RIGHT_M_IN)
    def x_for(col, cols):
        total = cols*box_w + max(0, cols-1)*x_gap
        start = Inches(LEFT_M_IN) + max((page_width - total) / 2, 0)
        return start + col*(box_w + x_gap)

    centers = {}
    row_index_by_id = {}
    col_index_by_id = {}

    current_top_y = Inches(TOP_M_IN) + Inches(TITLE_H_IN)
    for row_idx, (is_ext, nodes) in enumerate(layout_rows):
        cols = len(nodes)
        for col, node_id in enumerate(nodes):
            if is_ext:
                mgr_row = by_id_all[node_id]
                x = x_for(col, cols)
                draw_person_box(slide, x, current_top_y, box_w, box_h,
                                getattr(mgr_row, "name", ""), getattr(mgr_row, "title", ""), getattr(mgr_row, "tenure", ""),
                                name_pt, title_pt, tenure_pt)
                centers[node_id] = (x + box_w/2, current_top_y + box_h/2)
            else:
                r = by_id[node_id]
                x = x_for(col, cols)
                draw_person_box(slide, x, current_top_y, box_w, box_h,
                                getattr(r, "name", ""), getattr(r, "title", ""), getattr(r, "tenure", ""),
                                name_pt, title_pt, tenure_pt)
                centers[node_id] = (x + box_w/2, current_top_y + box_h/2)
            row_index_by_id[node_id] = row_idx
            col_index_by_id[node_id] = col
        current_top_y += (box_h + y_gap)

    thin = Emu(3)
    GAP_SAFE = Emu(10)
    JITTER_STEP = Emu(8)
    JITTER_BUCKETS = 5

    def safe_bus_y(m_bottom_y, s_top_y):
        top = m_bottom_y + GAP_SAFE
        bottom = s_top_y - GAP_SAFE
        if bottom <= top:
            return (m_bottom_y + s_top_y) / 2
        return (top + bottom) / 2

    def bus_y_for(parent_id, m_bottom_y, s_top_y):
        base = safe_bus_y(m_bottom_y, s_top_y)
        seed = (row_index_by_id.get(parent_id, 0) * 131 + col_index_by_id.get(parent_id, 0) * 17 + hash(parent_id)) & 0xFFFFFFFF
        lane = (seed % JITTER_BUCKETS) - (JITTER_BUCKETS // 2)
        y = base + lane * JITTER_STEP
        top = m_bottom_y + GAP_SAFE
        bottom = s_top_y - GAP_SAFE
        return max(top, min(bottom, y))

    def draw_bus_connect(mx, m_bottom_y, sx, s_top_y, parent_id):
        by = bus_y_for(parent_id, m_bottom_y, s_top_y)
        if abs(by - m_bottom_y) > 0:
            draw_line(slide, mx - thin/2, m_bottom_y, thin, by - m_bottom_y)
        if abs(sx - mx) > 0:
            draw_line(slide, min(mx, sx), by - thin/2, abs(sx - mx), thin)
        if abs(s_top_y - by) > 0:
            draw_line(slide, sx - thin/2, by, thin, s_top_y - by)

    for mng, subs in children.items():
        if mng not in centers:
            continue
        mx, my = centers[mng]
        m_bottom_y = my + box_h/2
        for s in subs:
            if s not in centers:
                continue
            sx, sy = centers[s]
            s_top_y = sy - box_h/2
            draw_bus_connect(mx, m_bottom_y, sx, s_top_y, mng)

    for mid in ext_mgrs_order:
        if mid not in centers:
            continue
        mx, my = centers[mid]
        m_bottom_y = my + box_h/2
        for child in mgr_to_dept_roots.get(mid, []):
            if child not in centers:
                continue
            sx, sy = centers[child]
            s_top_y = sy - box_h/2
            draw_bus_connect(mx, m_bottom_y, sx, s_top_y, mid)

    add_footer_pagenum(slide, prs, page_num, total_pages)

# ---------------------------
# PPT Builder
# ---------------------------

def build_ppt(df: pd.DataFrame) -> BytesIO:
    for c in REQUIRED_CANON:
        if c not in df.columns:
            raise ValueError(f"Missing column: {c}")
    df = df.fillna("")
    df = df.astype({c: str for c in REQUIRED_CANON})
    df = df.drop_duplicates(subset=["employee_id"])

    by_id_all = {r.employee_id: r for r in df.itertuples(index=False)}

    prs = Presentation()
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)

    dept_ordered = get_department_order(df)
    dept_counts = df.groupby("department")["employee_id"].nunique().to_dict()
    total_employees = df["employee_id"].nunique()

    total_pages = 1 + len(dept_ordered)
    current_page = 1

    draw_cover_slide(prs, dept_ordered, dept_counts, total_employees,
                     page_num=current_page, total_pages=total_pages, df=df)
    current_page += 1

    for dept in dept_ordered:
        sub = df[df["department"] == dept]
        if sub.empty:
            continue
        draw_department_slide(prs, dept_name=dept, sub_df=sub,
                              page_num=current_page, total_pages=total_pages,
                              by_id_all=by_id_all)
        current_page += 1

    bio = BytesIO()
    prs.save(bio); bio.seek(0)
    return bio

# ---------------------------
# Column mapping helpers
# ---------------------------

def _norm(s: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", str(s).lower()).strip()

def guess_mapping(columns):
    mapping = {}
    normcols = [(c, _norm(c)) for c in columns]
    for canon, keys in SYNONYMS.items():
        keys_norm = [_norm(k) for k in keys]
        best = None
        # exact
        for orig, norm in normcols:
            if norm in keys_norm or norm == canon:
                best = orig; break
        # startswith
        if not best:
            for orig, norm in normcols:
                if any(norm.startswith(k) for k in keys_norm):
                    best = orig; break
        # contains
        if not best:
            for orig, norm in normcols:
                if any(k in norm for k in keys_norm):
                    best = orig; break
        if best:
            mapping[canon] = best
    return mapping

def apply_mapping(df: pd.DataFrame, mapping: dict) -> pd.DataFrame:
    out = df.copy()
    # rename selected columns to canonical names
    rename_map = {src: canon for canon, src in mapping.items() if src in out.columns and src}
    out = out.rename(columns=rename_map)
    # ensure all canonical columns exist
    for c in ALL_CANON:
        if c not in out.columns:
            out[c] = ""
    return out

def render_mapping_ui(df: pd.DataFrame) -> Optional[pd.DataFrame]:
    st.subheader("Map your columns")
    st.caption("Choose which uploaded columns map to the org-chart fields (required fields marked with *).")

    columns = list(df.columns)
    guessed = guess_mapping(columns)

    if "mapping" not in st.session_state:
        st.session_state.mapping = {k: guessed.get(k, "") for k in ALL_CANON}

    with st.form("mapping_form", clear_on_submit=False):
        cols = st.columns(3)
        for i, canon in enumerate(ALL_CANON):
            label = f"{canon}{' *' if canon in REQUIRED_CANON else ''}"
            with cols[i % 3]:
                default_idx = 0
                if canon in guessed and guessed[canon] in columns:
                    default_idx = columns.index(guessed[canon]) + 1
                st.session_state.mapping[canon] = st.selectbox(
                    label,
                    options=["(not mapped)"] + columns,
                    index=default_idx,
                    key=f"map_{canon}",
                )

        dupes = []
        chosen = [v for v in st.session_state.mapping.values() if v and v != "(not mapped)"]
        for val in set(chosen):
            if chosen.count(val) > 1:
                dupes.append(val)
        if dupes:
            st.warning(f"Same source column mapped multiple times: {', '.join(sorted(set(dupes)))}")

        submitted = st.form_submit_button("Apply mapping")

    if submitted:
        clean_map = {
            canon: (src if src and src != "(not mapped)" else "")
            for canon, src in st.session_state.mapping.items()
        }
        missing_required = [c for c in REQUIRED_CANON if not clean_map.get(c)]
        if missing_required:
            st.error(f"Please map required fields: {', '.join(missing_required)}")
            return None

        mapped = apply_mapping(df, clean_map)

        # Fill friendly tenure text if tenure empty but tenure_calc present
        if "tenure" in mapped.columns and "tenure_calc" in mapped.columns:
            mask = (mapped["tenure"].astype(str).str.strip() == "") & (mapped["tenure_calc"].astype(str).str.strip() != "")
            mapped.loc[mask, "tenure"] = mapped.loc[mask, "tenure_calc"].astype(str).str.strip() + " yrs"

        st.success("Mapping applied.")
        return mapped

    return None

def auto_map_or_passthrough(df: pd.DataFrame) -> Optional[pd.DataFrame]:
    """
    If AUTO_MAP_WHEN_DISABLED is True:
      - Guess mapping from uploaded headers to canonical fields.
      - If all REQUIRED_CANON are guessed, apply and return the mapped df.
      - Else show an error prompting to enable mapping UI.
    If AUTO_MAP_WHEN_DISABLED is False:
      - Assume headers are already canonical; validate and return or error.
    """
    if AUTO_MAP_WHEN_DISABLED:
        guessed = guess_mapping(df.columns)
        missing = [c for c in REQUIRED_CANON if c not in guessed]
        if missing:
            st.error(
                "Required fields could not be auto-mapped: "
                + ", ".join(missing)
                + ". Set ENABLE_MAPPING_UI = True to map manually."
            )
            return None
        mapped = apply_mapping(df, guessed)
        with st.expander("Auto-mapping details", expanded=False):
            st.write({k: guessed.get(k, None) for k in ALL_CANON})
        st.success("Auto-mapping applied.")
        # tenure text convenience
        if "tenure" in mapped.columns and "tenure_calc" in mapped.columns:
            mask = (mapped["tenure"].astype(str).str.strip() == "") & (mapped["tenure_calc"].astype(str).str.strip() != "")
            mapped.loc[mask, "tenure"] = mapped.loc[mask, "tenure_calc"].astype(str).str.strip() + " yrs"
        return mapped
    else:
        missing = [c for c in REQUIRED_CANON if c not in df.columns]
        if missing:
            st.error(
                "Missing required headers: "
                + ", ".join(missing)
                + ". Either rename your columns or set ENABLE_MAPPING_UI = True."
            )
            return None
        return df

# ---------------------------
# Streamlit UI
# ---------------------------

uploaded = st.file_uploader("Upload CSV or Excel", type=["csv", "xlsx", "xls"])

raw_df: Optional[pd.DataFrame] = None
if uploaded is not None:
    try:
        name = (uploaded.name or "").lower()
        if name.endswith(".xlsx") or name.endswith(".xls"):
            try:
                raw_df = pd.read_excel(uploaded, dtype=str)
            except Exception:
                if name.endswith(".xlsx"):
                    raw_df = pd.read_excel(uploaded, dtype=str, engine="openpyxl")
                else:
                    raw_df = pd.read_excel(uploaded, dtype=str, engine="xlrd")
        else:
            raw_df = pd.read_csv(uploaded, dtype=str)
        raw_df = raw_df.fillna("")
        st.success(f"Loaded file with {len(raw_df):,} rows and {len(raw_df.columns)} columns.")
        st.dataframe(raw_df.head(20), use_container_width=True)
    except Exception as e:
        st.error(f"Unable to read file: {e}")
        raw_df = None
else:
    st.info("Upload a file to begin.")

mapped_df: Optional[pd.DataFrame] = None
if raw_df is not None:
    if ENABLE_MAPPING_UI:
        mapped_df = render_mapping_ui(raw_df)
    else:
        mapped_df = auto_map_or_passthrough(raw_df)

# ---- Editor + Summary ----
edited_df: Optional[pd.DataFrame] = None
if mapped_df is not None:
    st.subheader("Edit data (after mapping)")
    st.caption("Click cells to edit, press Enter to commit. Use + to add rows.")

    edited_df = DATA_EDITOR(
        mapped_df,
        num_rows="dynamic",
        use_container_width=True,
        hide_index=True,
        key="editor_main",
        column_config={
            "employee_id": st.column_config.TextColumn("employee_id", help="Unique ID", required=True),
            "name": st.column_config.TextColumn("name", required=True),
            "title": st.column_config.TextColumn("title", required=True),
            "manager_id": st.column_config.TextColumn("manager_id", help="Manager employee_id; blank for top-level"),
            "department": st.column_config.TextColumn("department", required=True),
            "tenure": st.column_config.TextColumn("tenure"),
            "tenure_calc": st.column_config.NumberColumn(
                "tenure_calc", help="Years as a number (e.g., 2.5)", step=0.1, format="%.1f"
            ),
        },
    )

    try:
        live = edited_df if edited_df is not None else mapped_df
        st.subheader("Department summary (live)")
        counts = live.groupby("department")["employee_id"].nunique().reset_index(name="Employees")
        if "tenure_calc" in live.columns:
            tmp = live.copy()
            tmp["tenure_calc"] = pd.to_numeric(tmp["tenure_calc"], errors="coerce")
            avg_df = tmp.groupby("department")["tenure_calc"].mean().round(1).reset_index(name="Avg Tenure (yrs)")
            counts = counts.merge(avg_df, on="department", how="left")
        st.dataframe(counts.sort_values("department"), use_container_width=True)
    except Exception as _e:
        st.caption(f"Summary note: {_e}")

# ---- Generate deck ----
col1, col2 = st.columns([3, 2])
with col1:
    out_name = st.text_input("Output filename", value="org_chart.pptx")
with col2:
    btn = st.button("Generate PowerPoint", type="primary")

if btn:
    use_df = edited_df if edited_df is not None else mapped_df
    if use_df is None:
        st.error("Please upload, (auto)map, and (optionally) edit your data first.")
    else:
        missing = [c for c in REQUIRED_CANON if c not in use_df.columns]
        if missing:
            st.error(f"Missing required columns: {', '.join(missing)}")
        else:
            try:
                ppt_bytes = build_ppt(use_df)
                st.success("Deck generated!")
                st.download_button(
                    label="Download PPTX",
                    data=ppt_bytes,
                    file_name=out_name or "org_chart.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            except Exception as e:
                st.error(f"Failed to generate deck: {e}")
