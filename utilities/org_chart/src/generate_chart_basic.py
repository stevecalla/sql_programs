# org_chart_basic.py
import pandas as pd
from collections import defaultdict, deque
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

INPUT = "org_chart.csv"
OUTPUT = "org_chart_basic.pptx"

df = pd.read_csv(INPUT, dtype=str).fillna("")
need = ["employee_id","name","title","manager_id"]
for c in need:
    if c not in df.columns:
        raise ValueError(f"Missing column: {c}")
    df[c] = df[c].astype(str).str.strip()
df = df.drop_duplicates(subset=["employee_id"])

by_id = {r.employee_id: r for r in df.itertuples(index=False)}
children = defaultdict(list)
for r in df.itertuples(index=False):
    if r.manager_id and r.manager_id in by_id and r.manager_id != r.employee_id:
        children[r.manager_id].append(r.employee_id)

roots = [
    r.employee_id for r in df.itertuples(index=False)
    if (r.manager_id == "" or r.manager_id not in by_id or r.manager_id == r.employee_id)
]
if not roots:
    all_children = {c for subs in children.values() for c in subs}
    roots = [eid for eid in by_id if eid not in all_children] or [next(iter(by_id))]

# BFS levels
level = {}
for rt in roots:
    q = deque([(rt, 0)])
    while q:
        u, d = q.popleft()
        if u in level and level[u] <= d:
            continue
        level[u] = d
        for v in children[u]:
            q.append((v, d + 1))

# stable DFS per level
ordered_by_level = defaultdict(list)
def order_subtree(u):
    ordered_by_level[level[u]].append(u)
    for v in children[u]:
        order_subtree(v)
for r in roots:
    order_subtree(r)

# draw (NO CONNECTORS)
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

left_margin, top_margin = Inches(0.5), Inches(0.6)
x_gap, y_gap = Inches(2.6), Inches(1.6)
box_w, box_h = Inches(2.4), Inches(1.0)

page_width = prs.slide_width - 2 * left_margin

def x_for(col, cols):
    total = cols * x_gap
    start = left_margin + max((page_width - total) / 2, 0)
    return start + col * x_gap

for lvl in sorted(ordered_by_level):
    row = ordered_by_level[lvl]
    for col, emp_id in enumerate(row):
        r = by_id[emp_id]
        x = x_for(col, len(row))
        y = top_margin + lvl * y_gap

        box = slide.shapes.add_shape(
            autoshape_type_id=1,  # rectangle
            left=x, top=y, width=box_w, height=box_h
        )
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(240,240,240)
        box.line.color.rgb = RGBColor(120,120,120)

        tf = box.text_frame
        tf.text = r.name
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        p2 = tf.add_paragraph()
        p2.text = r.title
        p2.font.size = Pt(10)

prs.save(OUTPUT)
print(f"Saved {OUTPUT}")
