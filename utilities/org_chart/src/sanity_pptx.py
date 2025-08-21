# sanity_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
slide.shapes.title.text = "Sanity Check"
slide.placeholders[1].text = "If this opens, python-pptx is working."

out = "sanity_ok.pptx"
prs.save(out)
print(f"Saved {out}")
